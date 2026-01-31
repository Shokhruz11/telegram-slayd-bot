# -*- coding: utf-8 -*-
"""
SUPER TALABA PRO BOT (Gemini)

Funksiyalar:
- AI bilan suhbat (Gemini)
- Slayd (PPTX): 📝 Slayd
- Mustaqil ish / referat (DOCX): 📄 Mustaqil ish / Referat
- Kurs ishi (DOCX): 📚 Kurs ishi
- Insho (DOCX): ✍️ Insho
- Tezis (DOCX): 📌 Tezislar
- Maqola (DOCX): 📰 Maqola
- To'lov + chek tizimi (skrinshot → admin → tasdiq → fayl)
- Referal tizimi (referal link, bonus buyurtmalar)
- Profi jamoa bo'limi (admin bilan aloqaga yo'naltirish)
"""

import os
import sqlite3
from datetime import datetime
import random
import string

import telebot
from telebot import types

import google.generativeai as genai

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from docx import Document

# ============================
#       ENV SOZLAMALAR
# ============================

BOT_TOKEN = os.getenv("BOT_TOKEN", "")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY", "")

ADMIN_TELEGRAM_ID = os.getenv("ADMIN_TELEGRAM_ID", "")
try:
    ADMIN_TELEGRAM_ID = int(ADMIN_TELEGRAM_ID) if ADMIN_TELEGRAM_ID else 0
except Exception:
    ADMIN_TELEGRAM_ID = 0

ADMIN_USERNAME = os.getenv("ADMIN_USERNAME", "@admin")

CARD_NUMBER = os.getenv("CARD_NUMBER", "0000 0000 0000 0000")
CARD_OWNER = os.getenv("CARD_OWNER", "Karta egasi")

try:
    MAX_FREE_LIMIT = int(os.getenv("MAX_FREE_LIMIT", "1"))
except Exception:
    MAX_FREE_LIMIT = 1

try:
    REF_BONUS_PER_PAID = int(os.getenv("REF_BONUS_PER_PAID", "1"))
except Exception:
    REF_BONUS_PER_PAID = 1

PRICE_PER_ORDER = 5000       # so'm, 20 betgacha
MAX_PAGES_PER_ORDER = 20

DB_NAME = "super_talaba.sqlite3"
FILES_DIR = "generated_files"

if not BOT_TOKEN:
    print("❌ BOT_TOKEN topilmadi!")
if not GEMINI_API_KEY:
    print("❌ GEMINI_API_KEY topilmadi!")

bot = telebot.TeleBot(BOT_TOKEN) if BOT_TOKEN else None

# Gemini sozlash
if GEMINI_API_KEY:
    try:
        genai.configure(api_key=GEMINI_API_KEY)
        gemini_model = genai.GenerativeModel("gemini-1.5-flash")
        print("✅ Gemini tayyor.")
    except Exception as e:
        print("❌ Gemini xatosi:", e)
        gemini_model = None
else:
    gemini_model = None

# xotiradagi state
user_state = {}   # chat_id -> state string
user_data = {}    # chat_id -> dict


# ============================
#    BAZA VA FUNKSIYALAR
# ============================

def ensure_files_dir():
    if not os.path.exists(FILES_DIR):
        os.makedirs(FILES_DIR)


def init_db():
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()

    c.execute("""
    CREATE TABLE IF NOT EXISTS users (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        telegram_id INTEGER UNIQUE,
        username TEXT,
        full_name TEXT,
        free_used INTEGER DEFAULT 0,
        bonus_orders INTEGER DEFAULT 0,
        ref_code TEXT UNIQUE,
        invited_by INTEGER,
        created_at TEXT
    )
    """)

    c.execute("""
    CREATE TABLE IF NOT EXISTS orders (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER,
        work_type TEXT,
        topic TEXT,
        pages INTEGER,
        price INTEGER,
        status TEXT,
        file_path TEXT,
        created_at TEXT
    )
    """)

    c.execute("""
    CREATE TABLE IF NOT EXISTS payments (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        order_id INTEGER,
        user_id INTEGER,
        amount INTEGER,
        status TEXT,
        screenshot_file_id TEXT,
        created_at TEXT
    )
    """)

    conn.commit()
    conn.close()


def generate_ref_code(length: int = 6) -> str:
    chars = string.ascii_uppercase + string.digits
    return ''.join(random.choice(chars) for _ in range(length))


def get_or_create_user(message, ref_payload: str | None = None):
    """
    Foydalanuvchini bazaga yozadi.
    Agar start'da referal kod kelsa, invited_by ni belgilaydi.
    """
    tg_id = message.from_user.id
    username = message.from_user.username or ""
    full_name = (message.from_user.first_name or "") + " " + (message.from_user.last_name or "")

    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("SELECT id, free_used, bonus_orders, ref_code, invited_by FROM users WHERE telegram_id = ?", (tg_id,))
    row = c.fetchone()

    if row:
        user_id, free_used, bonus_orders, ref_code, invited_by = row
        conn.close()
        return user_id
    else:
        # yangi foydalanuvchi
        now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        # referal kodi yaratamiz
        ref_code = generate_ref_code()
        invited_by_id = None

        # agar start'dan ref_XXXX kelsa, taklif qilgan odamni topamiz
        if ref_payload and ref_payload.startswith("ref_"):
            code = ref_payload[4:].strip().upper()
            c.execute("SELECT id FROM users WHERE ref_code = ?", (code,))
            ref_row = c.fetchone()
            if ref_row:
                invited_by_id = ref_row[0]

        c.execute("""
            INSERT INTO users (telegram_id, username, full_name, free_used, bonus_orders, ref_code, invited_by, created_at)
            VALUES (?, ?, ?, 0, 0, ?, ?, ?)
        """, (tg_id, username, full_name, ref_code, invited_by_id, now))
        conn.commit()
        user_id = c.lastrowid
        conn.close()
        return user_id


def get_user_by_id(user_id: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("""
        SELECT id, telegram_id, username, full_name, free_used, bonus_orders, ref_code, invited_by
        FROM users WHERE id = ?
    """, (user_id,))
    row = c.fetchone()
    conn.close()
    return row


def get_user_by_telegram_id(tg_id: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("""
        SELECT id, telegram_id, username, full_name, free_used, bonus_orders, ref_code, invited_by
        FROM users WHERE telegram_id = ?
    """, (tg_id,))
    row = c.fetchone()
    conn.close()
    return row


def update_user_free_used(user_id: int, delta: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("UPDATE users SET free_used = free_used + ? WHERE id = ?", (delta, user_id))
    conn.commit()
    conn.close()


def update_user_bonus_orders(user_id: int, delta: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("UPDATE users SET bonus_orders = MAX(bonus_orders + ?, 0) WHERE id = ?", (delta, user_id))
    conn.commit()
    conn.close()


def get_user_balance_info(user_id: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("SELECT free_used, bonus_orders, ref_code FROM users WHERE id = ?", (user_id,))
    row = c.fetchone()

    c.execute("SELECT COUNT(*) FROM orders WHERE user_id = ? AND status = 'done'", (user_id,))
    done_count = c.fetchone()[0]
    conn.close()

    if row:
        free_used, bonus_orders, ref_code = row
        return int(free_used), int(bonus_orders), ref_code, done_count
    else:
        return 0, 0, None, 0


def create_order(user_id: int, work_type: str, topic: str, pages: int, price: int, status: str):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    c.execute("""
        INSERT INTO orders (user_id, work_type, topic, pages, price, status, file_path, created_at)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?)
    """, (user_id, work_type, topic, pages, price, status, "", now))
    conn.commit()
    order_id = c.lastrowid
    conn.close()
    return order_id


def get_order(order_id: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("""
        SELECT id, user_id, work_type, topic, pages, price, status, file_path
        FROM orders WHERE id = ?
    """, (order_id,))
    row = c.fetchone()
    conn.close()
    return row


def update_order_status(order_id: int, status: str, file_path: str | None = None):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    if file_path is not None:
        c.execute("UPDATE orders SET status = ?, file_path = ? WHERE id = ?", (status, file_path, order_id))
    else:
        c.execute("UPDATE orders SET status = ? WHERE id = ?", (status, order_id))
    conn.commit()
    conn.close()


def create_payment(order_id: int, user_id: int, amount: int, screenshot_file_id: str):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    c.execute("""
        INSERT INTO payments (order_id, user_id, amount, status, screenshot_file_id, created_at)
        VALUES (?, ?, ?, ?, ?, ?)
    """, (order_id, user_id, amount, "pending", screenshot_file_id, now))
    conn.commit()
    payment_id = c.lastrowid
    conn.close()
    return payment_id


def get_payment(payment_id: int):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("""
        SELECT id, order_id, user_id, amount, status, screenshot_file_id
        FROM payments WHERE id = ?
    """, (payment_id,))
    row = c.fetchone()
    conn.close()
    return row


def update_payment_status(payment_id: int, status: str):
    conn = sqlite3.connect(DB_NAME)
    c = conn.cursor()
    c.execute("UPDATE payments SET status = ? WHERE id = ?", (status, payment_id))
    conn.commit()
    conn.close()


# ============================
#        GEMINI FUNKSIYA
# ============================

def ask_ai(prompt: str) -> str:
    if not gemini_model:
        return "ERROR: AI sozlanmagan. Admin GEMINI_API_KEY ni tekshirishi kerak."
    try:
        resp = gemini_model.generate_content(prompt)
        text = getattr(resp, "text", "") or ""
        text = text.strip()
        if not text:
            return "ERROR: AI javob qaytara olmadi."
        return text
    except Exception as e:
        print("Gemini xatosi:", repr(e))
        return "ERROR: AI xizmatida xatolik yuz berdi."


# ============================
#       FAYL YARATISH
# ============================

def create_pptx_from_text(topic: str, ai_text: str) -> str:
    ensure_files_dir()
    prs = Presentation()

    colors = [
        RGBColor(0, 102, 204),
        RGBColor(34, 139, 34),
        RGBColor(128, 0, 128),
        RGBColor(220, 20, 60),
        RGBColor(255, 140, 0),
        RGBColor(47, 79, 79),
    ]

    blocks = [b.strip() for b in ai_text.split("\n\n") if b.strip()]
    if not blocks:
        blocks = [ai_text]

    for i, block in enumerate(blocks):
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)

        bg = slide.shapes.add_shape(
            autoshape_type_id=1,
            left=Inches(0),
            top=Inches(0),
            width=Inches(13.3),
            height=Inches(7.5),
        )
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = colors[i % len(colors)]
        bg.line.width = 0

        tx = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(1),
            width=Inches(12),
            height=Inches(5.5),
        )
        tf = tx.text_frame
        tf.word_wrap = True

        lines = block.split("\n")
        title = lines[0][:80]
        body = "\n".join(lines[1:]) if len(lines) > 1 else ""

        p_title = tf.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(36)
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(255, 255, 255)
        p_title.alignment = PP_ALIGN.CENTER

        if body:
            p_body = tf.add_paragraph()
            p_body.text = body
            p_body.font.size = Pt(24)
            p_body.font.color.rgb = RGBColor(255, 255, 255)
            p_body.alignment = PP_ALIGN.LEFT

    filename = os.path.join(FILES_DIR, f"slayd_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    prs.save(filename)
    return filename


def create_docx_from_text(title: str, ai_text: str, work_type: str) -> str:
    ensure_files_dir()
    doc = Document()
    doc.add_heading(f"{work_type}: {title}", level=1)

    blocks = [b.strip() for b in ai_text.split("\n\n") if b.strip()]
    if not blocks:
        blocks = [ai_text]

    for block in blocks:
        doc.add_paragraph(block)

    filename = os.path.join(FILES_DIR, f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx")
    doc.save(filename)
    return filename


def generate_order_file(order_id: int) -> bool:
    order = get_order(order_id)
    if not order:
        return False

    o_id, user_id, work_type, topic, pages, price, status, file_path = order

    user = get_user_by_id(user_id)
    if not user:
        return False
    _, tg_id, _, _, _, _, _, _ = user
    chat_id = tg_id

    if work_type == "slayd":
        prompt = (
            f"Mavzu: {topic}\n"
            f"{pages} ta slayd uchun matn tuzing. Har bir slayd uchun sarlavha va 3–6 ta punkt yozing.\n"
            "Har bir slayd matnini ikki bo'sh qator bilan ajrating."
        )
        answer = ask_ai(prompt)
        if answer.startswith("ERROR:"):
            bot.send_message(chat_id, answer[6:])
            return False

        filename = create_pptx_from_text(topic, answer)
        update_order_status(order_id, "done", filename)
        with open(filename, "rb") as f:
            bot.send_document(chat_id, f, caption=f"Slayd tayyor!\nMavzu: {topic}")
        return True

    else:
        if work_type == "referat":
            work_name = "Mustaqil ish / Referat"
        elif work_type == "kurs":
            work_name = "Kurs ishi"
        elif work_type == "insho":
            work_name = "Insho"
        elif work_type == "tezis":
            work_name = "Tezislar"
        else:
            work_name = "Maqola"

        prompt = (
            f"Mavzu: {topic}\n"
            f"Taxminan {pages} betlik {work_name} matnini yozing.\n"
            "Kirish, asosiy qism va xulosani alohida bo'limlarda, ilmiy va tushunarli uslubda bayon qiling."
        )
        if work_type == "tezis":
            prompt += "\nMatnni asosiy tezislar ko'rinishida, punktlar bilan yozing."

        answer = ask_ai(prompt)
        if answer.startswith("ERROR:"):
            bot.send_message(chat_id, answer[6:])
            return False

        filename = create_docx_from_text(topic, answer, work_name)
        update_order_status(order_id, "done", filename)
        with open(filename, "rb") as f:
            bot.send_document(chat_id, f, caption=f"{work_name} tayyor!\nMavzu: {topic}")
        return True


# ============================
#           MENYU
# ============================

def main_menu():
    kb = types.ReplyKeyboardMarkup(resize_keyboard=True)
    kb.row("📝 Slayd", "📄 Mustaqil ish / Referat")
    kb.row("📚 Kurs ishi", "✍️ Insho", "📌 Tezislar")
    kb.row("📰 Maqola", "🤖 AI bilan suhbat")
    kb.row("🧑‍🏫 Profi jamoa", "🎁 Referal bonus")
    kb.row("💰 Balans", "💵 To'lov / Chek", "❓ Yordam")
    return kb


# ============================
#         KOMANDALAR
# ============================

@bot.message_handler(commands=["start"])
def cmd_start(message):
    # referal payload: /start ref_XXXX
    ref_payload = None
    if " " in message.text:
        parts = message.text.split(" ", 1)
        ref_payload = parts[1].strip()

    user_id = get_or_create_user(message, ref_payload)
    free_used, bonus_orders, ref_code, done_count = get_user_balance_info(user_id)

    text = (
        "Assalomu alaykum! 😊\n\n"
        "Siz *SUPER TALABA PRO BOT*dasiz.\n\n"
        "Bu bot orqali:\n"
        "• 📝 Slayd (PPTX)\n"
        "• 📄 Mustaqil ish / referat\n"
        "• 📚 Kurs ishi\n"
        "• ✍️ Insho\n"
        "• 📌 Tezis\n"
        "• 📰 Maqola\n"
        "ni avtomatik tayyorlab olishingiz mumkin.\n\n"
        f"🎁 Bepul buyurtma limiti: {MAX_FREE_LIMIT} ta\n"
        f"✅ Siz foydalanilgan: {free_used} ta\n"
        f"⭐ Referal bonus buyurtmalar: {bonus_orders} ta\n"
        f"📂 Yakunlangan buyurtmalar: {done_count} ta\n\n"
        "Menyudan kerakli bo'limni tanlang."
    )
    bot.send_message(message.chat.id, text, reply_markup=main_menu(), parse_mode="Markdown")


@bot.message_handler(commands=["help"])
def cmd_help(message):
    text = (
        "❓ *Yordam bo'limi*\n\n"
        "1) Xizmat turi bo'yicha tugmani bosing (Slayd, Referat, Kurs ishi...)\n"
        "2) Mavzuni kiriting.\n"
        "3) Bet / slayd sonini kiriting.\n"
        "4) Agar sizda bepul limiti bo'lsa, buyurtma avtomatik bajariladi.\n"
        "5) Aks holda karta raqamiga to'lov qilib, chekni skrinshotini yuborasiz.\n"
        "6) Admin tasdiqlaydi → fayl tayyor bo'ladi.\n\n"
        "Referal bo'limida do'stlaringizni taklif qilib, qo'shimcha bepul buyurtmalar olishingiz mumkin."
    )
    bot.send_message(message.chat.id, text, reply_markup=main_menu(), parse_mode="Markdown")


@bot.message_handler(commands=["ai"])
def cmd_ai(message):
    user_state[message.chat.id] = "ai_chat"
    bot.send_message(message.chat.id, "Savolingizni yoki mavzuni yozing. AI javob qaytaradi.")


# ============================
#      MENYU BOSIMLARI
# ============================

@bot.message_handler(func=lambda m: m.text == "🤖 AI bilan suhbat")
def handle_ai_button(message):
    cmd_ai(message)


@bot.message_handler(func=lambda m: m.text == "❓ Yordam")
def handle_help_btn(message):
    cmd_help(message)


@bot.message_handler(func=lambda m: m.text == "🧑‍🏫 Profi jamoa")
def handle_pro_team(message):
    text = (
        "🧑‍🏫 *Profi jamoa*\n\n"
        "Murakkab kurs ishlari, diplom ishlari, magistratura uchun ishlarda "
        f"bevosita admin bilan bog'laning: {ADMIN_USERNAME}\n\n"
        "Bot orqali esa oddiy referat, kurs ishi, slayd va boshqa topshiriqlarni avtomatik bajartirishingiz mumkin."
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")


@bot.message_handler(func=lambda m: m.text == "🎁 Referal bonus")
def handle_referral(message):
    row = get_user_by_telegram_id(message.from_user.id)
    if not row:
        user_id = get_or_create_user(message)
        free_used, bonus_orders, ref_code, done_count = get_user_balance_info(user_id)
    else:
        user_id, tg_id, username, full_name, free_used, bonus_orders, ref_code, invited_by = row
        _, _, ref_code, done_count = get_user_balance_info(user_id)

    if not ref_code:
        ref_code = generate_ref_code()
        conn = sqlite3.connect(DB_NAME)
        c = conn.cursor()
        c.execute("UPDATE users SET ref_code = ? WHERE id = ?", (ref_code, user_id))
        conn.commit()
        conn.close()

    try:
        me = bot.get_me()
        bot_username = me.username
    except Exception:
        bot_username = "SuperTalabaBot"

    ref_link = f"https://t.me/{bot_username}?start=ref_{ref_code}"

    text = (
        "🎁 *Referal tizimi*\n\n"
        f"Sizning referal kodingiz: `{ref_code}`\n"
        f"Referal link: {ref_link}\n\n"
        f"Har bir do'stingiz ushbu link orqali kirib, pullik buyurtma qilsa,\n"
        f"sizga {REF_BONUS_PER_PAID} ta qo'shimcha bepul buyurtma qo'shiladi.\n\n"
        f"Joriy bonus buyurtmalar: {bonus_orders} ta."
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")


@bot.message_handler(func=lambda m: m.text == "💰 Balans")
def handle_balance(message):
    row = get_user_by_telegram_id(message.from_user.id)
    if not row:
        user_id = get_or_create_user(message)
    else:
        user_id = row[0]

    free_used, bonus_orders, ref_code, done_count = get_user_balance_info(user_id)
    text = (
        "💰 *Balans ma'lumotlari*\n\n"
        f"- Bepul limit: {MAX_FREE_LIMIT} ta\n"
        f"- Foydalanilgan bepul buyurtma: {free_used} ta\n"
        f"- Referal bonus buyurtmalar: {bonus_orders} ta\n"
        f"- Yakunlangan buyurtmalar: {done_count} ta\n\n"
        f"Har bir pullik buyurtma narxi: {PRICE_PER_ORDER} so'm (20 bet/slaydgacha)."
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")


@bot.message_handler(func=lambda m: m.text == "💵 To'lov / Chek")
def handle_payment_info(message):
    text = (
        "💵 *To'lov ma'lumotlari*\n\n"
        f"Karta raqami: `{CARD_NUMBER}`\n"
        f"Karta egasi: *{CARD_OWNER}*\n\n"
        "To'lovni amalga oshirgach, shu botga chek skrinshotini rasm ko'rinishida yuboring.\n"
        "Agar sizda faol buyurtma bo'lmasa, avval menyudan biror xizmat turini tanlab buyurtma yarating."
    )
    bot.send_message(message.chat.id, text, parse_mode="Markdown")


# --------- Xizmat tugmalari ---------

SERVICE_BUTTONS = {
    "📝 Slayd": "slayd",
    "📄 Mustaqil ish / Referat": "referat",
    "📚 Kurs ishi": "kurs",
    "✍️ Insho": "insho",
    "📌 Tezislar": "tezis",
    "📰 Maqola": "maqola",
}


@bot.message_handler(func=lambda m: m.text in SERVICE_BUTTONS.keys())
def handle_service_select(message):
    chat_id = message.chat.id
    work_type = SERVICE_BUTTONS[message.text]

    user_state[chat_id] = "enter_topic"
    user_data[chat_id] = {"work_type": work_type}

    if work_type == "slayd":
        prompt = "Slayd mavzusini yozing:"
    elif work_type == "referat":
        prompt = "Mustaqil ish / referat mavzusini yozing:"
    elif work_type == "kurs":
        prompt = "Kurs ishi mavzusini yozing:"
    elif work_type == "insho":
        prompt = "Insho mavzusini yozing:"
    elif work_type == "tezis":
        prompt = "Tezislar mavzusini yozing:"
    else:
        prompt = "Maqola mavzusini yozing:"

    bot.send_message(chat_id, prompt, reply_markup=types.ReplyKeyboardRemove())


@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "enter_topic")
def handle_enter_topic(message):
    chat_id = message.chat.id
    data = user_data.get(chat_id, {})
    work_type = data.get("work_type", "referat")

    topic = message.text.strip()
    data["topic"] = topic
    user_data[chat_id] = data

    if work_type == "slayd":
        text = "Nechta slayd kerak? (1–20):"
    else:
        text = "Taxminan necha bet kerak? (1–20):"

    user_state[chat_id] = "enter_pages"
    bot.send_message(chat_id, text)


@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "enter_pages")
def handle_enter_pages(message):
    chat_id = message.chat.id
    data = user_data.get(chat_id, {})
    work_type = data.get("work_type", "referat")
    topic = data.get("topic", "")

    try:
        pages = int(message.text.strip())
        if not 1 <= pages <= MAX_PAGES_PER_ORDER:
            raise ValueError()
    except ValueError:
        bot.send_message(chat_id, f"Iltimos, 1 dan {MAX_PAGES_PER_ORDER} gacha bo'lgan son kiriting.")
        return

    data["pages"] = pages
    user_data[chat_id] = data

    # foydalanuvchini olamiz
    user_row = get_user_by_telegram_id(message.from_user.id)
    if not user_row:
        user_id = get_or_create_user(message)
        free_used, bonus_orders, ref_code, done_count = get_user_balance_info(user_id)
    else:
        user_id = user_row[0]
        free_used, bonus_orders, ref_code, done_count = get_user_balance_info(user_id)

    # bepul yoki pullik?
    if free_used < MAX_FREE_LIMIT:
        # asosiy bepul limit
        bot.send_message(chat_id, "Sizda bepul buyurtma mavjud. Buyurtma bepul bajariladi. AI ishlayapti...")
        update_user_free_used(user_id, 1)
        order_id = create_order(user_id, work_type, topic, pages, 0, "processing")
        ok = generate_order_file(order_id)
        if ok:
            bot.send_message(chat_id, "Bepul buyurtmangiz yakunlandi ✅", reply_markup=main_menu())
        else:
            bot.send_message(chat_id, "Buyurtmani bajarishda xato yuz berdi.", reply_markup=main_menu())
        user_state[chat_id] = None
        user_data[chat_id] = {}
        return
    elif bonus_orders > 0:
        # referal bonus hisobidan
        bot.send_message(chat_id, "Sizda referal bonus buyurtma mavjud. Buyurtma bepul bajariladi. AI ishlayapti...")
        update_user_bonus_orders(user_id, -1)
        order_id = create_order(user_id, work_type, topic, pages, 0, "processing")
        ok = generate_order_file(order_id)
        if ok:
            bot.send_message(chat_id, "Bonus asosida bepul buyurtmangiz yakunlandi ✅", reply_markup=main_menu())
        else:
            bot.send_message(chat_id, "Buyurtmani bajarishda xato yuz berdi.", reply_markup=main_menu())
        user_state[chat_id] = None
        user_data[chat_id] = {}
        return
    else:
        # pullik rejim
        price = PRICE_PER_ORDER
        order_id = create_order(user_id, work_type, topic, pages, price, "pending_payment")
        data["order_id"] = order_id
        user_data[chat_id] = data
        user_state[chat_id] = "waiting_payment_screenshot"

        text = (
            "💳 *Pullik buyurtma*\n\n"
            f"- Turi: {work_type}\n"
            f"- Mavzu: {topic}\n"
            f"- Bet/slayd soni: {pages}\n"
            f"- Narx: {price} so'm\n\n"
            "Iltimos, quyidagi kartaga to'lovni amalga oshiring:\n"
            f"Karta raqami: `{CARD_NUMBER}`\n"
            f"Karta egasi: *{CARD_OWNER}*\n\n"
            "To'lov tugagach, shu chatga chek skrinshotini rasm ko'rinishida yuboring."
        )
        bot.send_message(chat_id, text, parse_mode="Markdown")


# ============================
#       CHEK / PHOTO
# ============================

@bot.message_handler(content_types=["photo"])
def handle_payment_photo(message):
    chat_id = message.chat.id
    state = user_state.get(chat_id)

    if state != "waiting_payment_screenshot":
        bot.send_message(chat_id, "Bu rasm to'lov cheki sifatida qabul qilinmadi. Avval buyurtma yarating.")
        return

    data = user_data.get(chat_id, {})
    order_id = data.get("order_id")
    if not order_id:
        bot.send_message(chat_id, "Buyurtma ma'lumotlari topilmadi. Qayta urinib ko'ring.")
        user_state[chat_id] = None
        return

    order = get_order(order_id)
    if not order:
        bot.send_message(chat_id, "Buyurtma topilmadi. Qayta urinib ko'ring.")
        user_state[chat_id] = None
        return

    o_id, user_id, work_type, topic, pages, price, status, file_path = order

    file_id = message.photo[-1].file_id
    payment_id = create_payment(order_id, user_id, price, file_id)

    bot.send_message(chat_id, "Chek qabul qilindi. Admin tasdiqlaganidan so'ng faylingiz tayyorlanadi.")

    # admin'ga yuborish
    if ADMIN_TELEGRAM_ID:
        caption = (
            "💵 Yangi to'lov cheki\n\n"
            f"Payment ID: {payment_id}\n"
            f"Order ID: {order_id}\n"
            f"User ID: {user_id}\n"
            f"Username: @{message.from_user.username or 'no_username'}\n"
            f"Ish turi: {work_type}\n"
            f"Mavzu: {topic}\n"
            f"Bet/slayd: {pages}\n"
            f"Summa: {price} so'm\n\n"
            "Tasdiqlash yoki rad etish uchun tugmalardan foydalaning."
        )
        kb = types.InlineKeyboardMarkup()
        kb.add(
            types.InlineKeyboardButton("✅ Tasdiqlash", callback_data=f"pay_ok_{payment_id}"),
            types.InlineKeyboardButton("❌ Rad etish", callback_data=f"pay_no_{payment_id}")
        )
        try:
            bot.send_photo(ADMIN_TELEGRAM_ID, file_id, caption=caption, reply_markup=kb)
        except Exception as e:
            print("Admin'ga chek yuborishda xato:", e)
    user_state[chat_id] = None
    user_data[chat_id] = {}


# ============================
#      ADMIN CALLBACK
# ============================

@bot.callback_query_handler(func=lambda c: c.data.startswith("pay_ok_") or c.data.startswith("pay_no_"))
def handle_payment_callback(call):
    if call.from_user.id != ADMIN_TELEGRAM_ID:
        bot.answer_callback_query(call.id, "Siz admin emassiz.")
        return

    data = call.data
    is_ok = data.startswith("pay_ok_")
    payment_id = int(data.split("_")[-1])

    payment = get_payment(payment_id)
    if not payment:
        bot.answer_callback_query(call.id, "To'lov topilmadi.")
        return

    p_id, order_id, user_id, amount, status, screenshot_file_id = payment
    order = get_order(order_id)
    if not order:
        bot.answer_callback_query(call.id, "Buyurtma topilmadi.")
        return

    o_id, u_id, work_type, topic, pages, price, o_status, file_path = order

    user = get_user_by_id(user_id)
    if not user:
        bot.answer_callback_query(call.id, "Foydalanuvchi topilmadi.")
        return

    _, tg_id, username, full_name, free_used, bonus_orders, ref_code, invited_by = user
    chat_id = tg_id

    if is_ok:
        update_payment_status(payment_id, "approved")
        update_order_status(order_id, "processing")
        bot.answer_callback_query(call.id, "To'lov tasdiqlandi. Fayl yaratilmoqda.")
        bot.send_message(chat_id, "To'lov tasdiqlandi. Faylingiz tayyorlanmoqda...")

        ok = generate_order_file(order_id)
        if ok:
            bot.send_message(chat_id, "Buyurtmangiz bajarildi ✅", reply_markup=main_menu())
            # referal bonus
            if invited_by and REF_BONUS_PER_PAID > 0:
                update_user_bonus_orders(invited_by, REF_BONUS_PER_PAID)
                inv_user = get_user_by_id(invited_by)
                if inv_user:
                    inv_tg_id = inv_user[1]
                    try:
                        bot.send_message(
                            inv_tg_id,
                            f"🎁 Sizning referalingiz buyurtma qildi va to'lov tasdiqlandi.\n"
                            f"Sizga {REF_BONUS_PER_PAID} ta bonus buyurtma qo'shildi."
                        )
                    except Exception:
                        pass
        else:
            bot.send_message(chat_id, "Fayl yaratishda xato yuz berdi.", reply_markup=main_menu())
    else:
        update_payment_status(payment_id, "rejected")
        update_order_status(order_id, "payment_rejected")
        bot.answer_callback_query(call.id, "To'lov rad etildi.")
        bot.send_message(
            chat_id,
            "Admin to'lovni tasdiqlamadi. Iltimos, chekingizni tekshirib, kerak bo'lsa qayta yuboring.",
            reply_markup=main_menu()
        )


# ============================
#        AI CHAT (/ai)
# ============================

@bot.message_handler(func=lambda m: user_state.get(m.chat.id) == "ai_chat")
def handle_ai_chat(message):
    chat_id = message.chat.id
    text = message.text.strip()
    bot.send_chat_action(chat_id, "typing")
    answer = ask_ai(text)
    if answer.startswith("ERROR:"):
        bot.send_message(chat_id, answer[6:], reply_markup=main_menu())
    else:
        bot.send_message(chat_id, answer, reply_markup=main_menu())
    user_state[chat_id] = None


# ============================
#      DEFAULT HANDLER
# ============================

@bot.message_handler(func=lambda m: True, content_types=["text"])
def fallback(message):
    if message.text.startswith("/"):
        bot.send_message(message.chat.id, "Noma'lum buyruq. /start yoki /help ni yozib ko'ring.")
    else:
        bot.send_message(
            message.chat.id,
            "Menyudan kerakli bo'limni tanlang.",
            reply_markup=main_menu()
        )


# ============================
#      BOTNI ISHGA TUSHIRISH
# ============================

if __name__ == "__main__":
    print("🚀 SUPER TALABA PRO BOT ishga tushmoqda...")
    init_db()
    ensure_files_dir()
    bot.infinity_polling(skip_pending=True)
